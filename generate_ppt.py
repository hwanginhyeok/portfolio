"""
포트폴리오 PPT 생성 스크립트
미니멀 다크 테마, 10장 구성
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 디자인 시스템 ──
BG_DARK = RGBColor(0x0C, 0x0E, 0x12)       # #0c0e12
BG_CARD = RGBColor(0x16, 0x18, 0x1E)       # #16181e
ACCENT = RGBColor(0x3B, 0x82, 0xF6)        # blue-500
ACCENT_DIM = RGBColor(0x1E, 0x40, 0xAF)    # blue-800
TEXT_PRIMARY = RGBColor(0xF1, 0xF5, 0xF9)   # slate-100
TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)  # slate-400
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)     # slate-500
GREEN = RGBColor(0x22, 0xC5, 0x5E)          # green-500
YELLOW = RGBColor(0xEA, 0xB3, 0x08)         # yellow-500
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER_COLOR = RGBColor(0x1E, 0x29, 0x3B)   # slate-800

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── 헬퍼 함수 ──
def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=14, color=TEXT_PRIMARY, bold=False,
             space_before=Pt(4), space_after=Pt(2), font_name='맑은 고딕',
             alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    return p

def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border:
        shape.line.color.rgb = BORDER_COLOR
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # rounded corner
    shape.adjustments[0] = 0.02
    return shape

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def add_tag(slide, left, top, text, color=ACCENT):
    w, h = Inches(1.8), Pt(28)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(
        min(255, color.red + 10) if color.red < 30 else color.red // 5,
        min(255, color.green + 10) if color.green < 30 else color.green // 5,
        min(255, color.blue + 10) if color.blue < 30 else color.blue // 5
    )
    shape.line.fill.background()
    shape.adjustments[0] = 0.3
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Consolas'
    p.alignment = PP_ALIGN.CENTER
    return shape


# ════════════════════════════════════════
# SLIDE 1: Cover
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# accent line top
add_rect(slide, 0, 0, SLIDE_W, Pt(4), ACCENT)

# name
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
             '황인혁', font_size=52, color=WHITE, bold=True)

# English name
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
             'Inhyeok Hwang', font_size=24, color=TEXT_SECONDARY,
             font_name='Consolas')

# tagline
add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.8),
             'Motor Control & Product Engineer', font_size=28, color=ACCENT,
             font_name='Consolas')

# one-liner
add_text_box(slide, Inches(1.5), Inches(4.8), Inches(9), Inches(1),
             '기계공학의 물리적 사고를 기반으로,\n전력제어 알고리즘을 개발하고 제품을 양산까지 끌고 가는 엔지니어.',
             font_size=16, color=TEXT_SECONDARY)

# contact bar at bottom
add_rect(slide, 0, Inches(6.5), SLIDE_W, Inches(1), RGBColor(0x0A, 0x0B, 0x0F))
tf = add_text_box(slide, Inches(1.5), Inches(6.7), Inches(10), Inches(0.5),
                  'github.com/hwanginhyeok  |  inhyeok.hwang@gintlab.com',
                  font_size=12, color=TEXT_MUTED, font_name='Consolas')


# ════════════════════════════════════════
# SLIDE 2: Profile & Career
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             'Profile', font_size=32, color=WHITE, bold=True, font_name='Consolas')
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

# Left: Profile info
left_x = Inches(0.8)
info_items = [
    ('이름', '황인혁 (Inhyeok Hwang)'),
    ('직책', 'GINT 전력제어개발팀 · 선임연구원'),
    ('역할', '전력제어 개발자 + 제품개발팀 주니어 PM'),
    ('학력', '건국대 기계공학부(학사) / 기계설계학과(석사)'),
    ('연구실', 'RBDO 연구실 — IPMSM 모델링 및 고장진단'),
    ('직무발명', '2건 출원 (초기위치추정, 극저온 기동)'),
]
y = Inches(1.3)
for label, value in info_items:
    tf = add_text_box(slide, left_x, y, Inches(5.5), Inches(0.35),
                      '', font_size=13)
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = f'{label}  '
    run_label.font.size = Pt(12)
    run_label.font.color.rgb = TEXT_MUTED
    run_label.font.name = '맑은 고딕'
    run_label.font.bold = True
    run_value = p.add_run()
    run_value.text = value
    run_value.font.size = Pt(13)
    run_value.font.color.rgb = TEXT_PRIMARY
    run_value.font.name = '맑은 고딕'
    y += Inches(0.38)

# Right: Career Timeline
right_x = Inches(7)
add_text_box(slide, right_x, Inches(1.0), Inches(5), Inches(0.5),
             'Career Timeline', font_size=18, color=WHITE, bold=True, font_name='Consolas')

timeline = [
    ('2025.01 ~ 현재', 'SS500 전력 시스템 + PM', '선임연구원',
     '48V 자율주행 방제기 · 전력제어 + 제품개발 PM\n이슈 37건+ · DVP 5건 · QC 5단계 구축'),
    ('2023.02 ~ 2024.12', 'EOP 모터제어 알고리즘', '주임연구원',
     '400W 센서리스 FOC · 직무발명 2건 출원\n정격 응답 0.56s · 극저온 기동 3배 단축'),
    ('2021.03 ~ 2023.02', 'PMSM 고장진단 연구', '석사과정',
     '120kW IPMSM Co-simulation\n현대자동차·현대엔지비 과제'),
]

y = Inches(1.6)
for period, title, pos, desc in timeline:
    # timeline dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, right_x, y + Pt(6), Pt(10), Pt(10))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()

    # vertical line
    if title != 'PMSM 고장진단 연구':
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       right_x + Pt(4), y + Pt(18), Pt(2), Inches(1.2))
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER_COLOR
        line.line.fill.background()

    tx = right_x + Inches(0.3)
    add_text_box(slide, tx, y - Pt(2), Inches(2.5), Inches(0.3),
                 period, font_size=10, color=TEXT_MUTED, font_name='Consolas')
    add_text_box(slide, tx, y + Pt(14), Inches(4), Inches(0.3),
                 f'{title}  |  {pos}', font_size=13, color=WHITE, bold=True)
    add_text_box(slide, tx, y + Pt(32), Inches(5), Inches(0.6),
                 desc, font_size=10, color=TEXT_SECONDARY)
    y += Inches(1.5)

# Bottom: positioning statement
add_rect(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.9), BG_CARD, border=True)
add_text_box(slide, Inches(1.2), Inches(6.35), Inches(10.5), Inches(0.6),
             '"소프트웨어가 물리 세계와 만나는 접점에서 제품을 만드는 사람"',
             font_size=15, color=ACCENT, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# SLIDE 3: Career Growth Trajectory
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             'Career Growth', font_size=32, color=WHITE, bold=True, font_name='Consolas')
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

# Three phase cards
phases = [
    ('대학원', '시뮬레이션', '2021 — 2023', [
        '120kW PMSM 모델링/해석',
        'Ansys Maxwell + MATLAB/Simulink',
        '논문 + 학회 (수상 2건)',
        '보조연구원',
    ], TEXT_MUTED),
    ('GINT — EOP', '실물 제어', '2023 — 2024', [
        '12V / 400W / 1모터',
        'NXP MCU + C',
        '직무발명 2건 출원',
        '알고리즘 개발자',
    ], ACCENT),
    ('GINT — SS500', '시스템 + PM', '2025 — 현재', [
        '48V / 3kW / 10+ 노드',
        'STM32 + CAN FD',
        '이슈 37건+ 관리',
        '개발자 + PM',
    ], GREEN),
]

card_w = Inches(3.5)
gap = Inches(0.4)
start_x = Inches(0.8)

for i, (phase, subtitle, period, items, accent) in enumerate(phases):
    x = start_x + i * (card_w + gap)
    y_card = Inches(1.5)

    # card bg
    add_rect(slide, x, y_card, card_w, Inches(4.2), BG_CARD, border=True)

    # accent top bar
    add_rect(slide, x, y_card, card_w, Pt(4), accent)

    # phase title
    add_text_box(slide, x + Inches(0.3), y_card + Inches(0.3), card_w - Inches(0.6), Inches(0.4),
                 phase, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y_card + Inches(0.7), card_w - Inches(0.6), Inches(0.3),
                 subtitle, font_size=16, color=accent, font_name='Consolas')
    add_text_box(slide, x + Inches(0.3), y_card + Inches(1.1), card_w - Inches(0.6), Inches(0.3),
                 period, font_size=11, color=TEXT_MUTED, font_name='Consolas')

    # items
    y_item = y_card + Inches(1.6)
    for item in items:
        add_text_box(slide, x + Inches(0.3), y_item, card_w - Inches(0.6), Inches(0.35),
                     f'→  {item}', font_size=12, color=TEXT_SECONDARY)
        y_item += Inches(0.38)

    # arrow between cards
    if i < 2:
        arrow_x = x + card_w + Inches(0.05)
        add_text_box(slide, arrow_x, Inches(3.2), Inches(0.3), Inches(0.5),
                     '→', font_size=28, color=TEXT_MUTED, font_name='Consolas',
                     alignment=PP_ALIGN.CENTER)

# Bottom message
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
             '시뮬레이션 → 실물 제어 → 시스템 + PM  |  기술적 깊이와 관리 역량을 동시에 키워온 궤적',
             font_size=14, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# SLIDE 4: EOP Project (1/2)
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             'Core Project 1 — 400W EOP 인버터', font_size=28, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.35),
             '현대트랜시스 EDS용 전동 오일펌프  |  2023 ~ 2024  |  모터제어 알고리즘 개발',
             font_size=13, color=TEXT_MUTED)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3))

# Problem 1
add_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.4), BG_CARD, border=True)
add_text_box(slide, Inches(1.1), Inches(1.75), Inches(5), Inches(0.3),
             '문제 1: 정지 상태에서 모터 위치를 모른다', font_size=15, color=YELLOW, bold=True)
add_text_box(slide, Inches(1.1), Inches(2.15), Inches(5), Inches(0.7),
             '센서리스 제어는 정지 시 회전자 위치 미상 → align 기동(강제 정렬)\n역회전 위험 + 느린 응답',
             font_size=11, color=TEXT_SECONDARY)
tf = add_text_box(slide, Inches(1.1), Inches(2.8), Inches(5), Inches(1),
                  '', font_size=12)
for line in ['→ Reluctance 특성 이용 초기위치추정 알고리즘 개발',
             '   초기 위치 파악: 0.44μs  |  역회전: 제거',
             '   정격 속도 응답: 요구 1.0s → 0.56s 달성',
             '   ★ 직무발명 출원']:
    color = GREEN if '★' in line else (ACCENT if '→' == line[0] else TEXT_PRIMARY)
    add_para(tf, line, font_size=12, color=color, bold=('★' in line))

# Problem 2
add_rect(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(2.4), BG_CARD, border=True)
add_text_box(slide, Inches(7.1), Inches(1.75), Inches(5), Inches(0.3),
             '문제 2: -40℃에서 모터가 탈조한다', font_size=15, color=YELLOW, bold=True)
add_text_box(slide, Inches(7.1), Inches(2.15), Inches(5), Inches(0.7),
             '극저온 → 윤활유 점도 상승 → 관로 저항 증가 → 불규칙 부하\n오픈루프 기동 시 탈조, 센서리스 전환까지 300초+',
             font_size=11, color=TEXT_SECONDARY)
tf = add_text_box(slide, Inches(7.1), Inches(2.8), Inches(5), Inches(1),
                  '', font_size=12)
for line in ['→ Eγ 값 기반 부하 실시간 추정 + 능동 속도 제어',
             '   기동 시간: 300s → 100s 이내 (3배 단축)',
             '   -35℃, 50rpm 안정 운전 달성',
             '   ★ 직무발명 출원']:
    color = GREEN if '★' in line else (ACCENT if '→' == line[0] else TEXT_PRIMARY)
    add_para(tf, line, font_size=12, color=color, bold=('★' in line))

# Problem 3
add_rect(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(1.5), BG_CARD, border=True)
add_text_box(slide, Inches(1.1), Inches(4.45), Inches(5), Inches(0.3),
             '문제 3: 고온에서 전력 효율 낮음', font_size=15, color=YELLOW, bold=True)
tf = add_text_box(slide, Inches(1.1), Inches(4.85), Inches(5), Inches(0.8),
                  '', font_size=12)
add_para(tf, '→ DPWM ↔ SVPWM 자동 절환 알고리즘', font_size=12, color=ACCENT)
add_para(tf, '   입력전력 1.0~3.8% 개선  |  FET 온도 1~6℃ 저감', font_size=12, color=TEXT_PRIMARY)

# Results summary
add_rect(slide, Inches(6.8), Inches(4.3), Inches(5.5), Inches(1.5), BG_CARD, border=True)
add_text_box(slide, Inches(7.1), Inches(4.45), Inches(5), Inches(0.3),
             '성과 요약', font_size=15, color=GREEN, bold=True)
tf = add_text_box(slide, Inches(7.1), Inches(4.85), Inches(5), Inches(0.8),
                  '', font_size=12)
add_para(tf, '직무발명 2건  |  응답 0.56s  |  기동 3배 단축', font_size=13, color=WHITE, bold=True)
add_para(tf, '400W → 150W(CEER) → HV 3개 바리에이션 확장', font_size=11, color=TEXT_SECONDARY)

# Tech stack bar
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.3),
             'NXP MC9S12ZVM  ·  FOC  ·  Sensorless (Extended EMF)  ·  SVPWM/DPWM  ·  FWC  ·  MTPA  ·  CAN 2.0B  ·  C',
             font_size=11, color=TEXT_MUTED, font_name='Consolas', alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# SLIDE 5: SS500 Project (1/2) — 전력제어
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             'Core Project 2 — SS500 Speed Sprayer', font_size=28, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.35),
             'PLUVA FRIEND GT-SS500  |  2025 ~ 현재  |  전력제어 + PM',
             font_size=13, color=TEXT_MUTED)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3))

# Product spec card
add_rect(slide, Inches(0.8), Inches(1.6), Inches(3.5), Inches(2.5), BG_CARD, border=True)
add_text_box(slide, Inches(1.1), Inches(1.75), Inches(3), Inches(0.3),
             '제품 스펙', font_size=14, color=ACCENT, bold=True)
specs = [
    '전원: 48V 460Ah LiFePO4',
    '구동: 3kW BLDC x2 (좌/우 독립)',
    '팬: Hobbywing X6 Plus (CAN FD)',
    'VCU: STM32F479VI',
    '통신: 49개 CAN 메시지, 10+ 노드',
    '자율주행: GPS RTK + IMU + 카메라 3종',
]
tf = add_text_box(slide, Inches(1.1), Inches(2.1), Inches(3), Inches(2), '', font_size=11)
for spec in specs:
    add_para(tf, spec, font_size=11, color=TEXT_SECONDARY)

# Root cause analysis card
add_rect(slide, Inches(4.6), Inches(1.6), Inches(7.5), Inches(2.5), BG_CARD, border=True)
add_text_box(slide, Inches(4.9), Inches(1.75), Inches(7), Inches(0.3),
             '문제점 도출 · 근본 원인 분석', font_size=14, color=YELLOW, bold=True)

cases = [
    ('Main Relay 반복 ON/OFF', '통신→FW→배선 순서로 추적', '솔레노이드 GND 바운스 (배선 설계)'),
    ('배터리 SOC 70.2% 부정확', 'BMS 캘리브→충방전 패턴 분석', 'ROM 로직 + 셀밸런싱 복합'),
    ('VCU FW 무한 반복', 'CanTimeoutCheck.c 코드리뷰', 'MR OFF 시 timeout 클리어 버그'),
]
y = Inches(2.15)
for symptom, process, cause in cases:
    tf = add_text_box(slide, Inches(4.9), y, Inches(7), Inches(0.5), '', font_size=11)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f'▸ {symptom}'
    run.font.size = Pt(11)
    run.font.color.rgb = TEXT_PRIMARY
    run.font.bold = True
    run.font.name = '맑은 고딕'
    add_para(tf, f'  판별: {process} → {cause}', font_size=10, color=TEXT_MUTED)
    y += Inches(0.55)

# Bottom cards: Algorithm + HW + Test
bottom_cards = [
    ('알고리즘 · 로직', [
        '제어권 상태머신 설계',
        'CAN timeout 대응 로직',
        '충전/방전 모드 개선',
        '솔레노이드 버그 수정',
    ]),
    ('HW 분석', [
        'VCU Rev01→02→03 변경 이력',
        '6층 PCB 거버 분석',
        '회로도 8장 분석',
        '샤시접지/노이즈 리서치',
    ]),
    ('시험 기획 · 검증', [
        'KGAT 분사 시험 15건 매트릭스',
        '펌프 20대 랜덤 샘플링',
        'IQC 검사 가이드 10건+',
        'FAN 사양 변경 검토',
    ]),
]

card_w = Inches(3.6)
for i, (title, items) in enumerate(bottom_cards):
    x = Inches(0.8) + i * (card_w + Inches(0.3))
    y_c = Inches(4.4)
    add_rect(slide, x, y_c, card_w, Inches(2.3), BG_CARD, border=True)
    add_text_box(slide, x + Inches(0.3), y_c + Inches(0.15), card_w - Inches(0.6), Inches(0.3),
                 title, font_size=13, color=ACCENT, bold=True)
    tf = add_text_box(slide, x + Inches(0.3), y_c + Inches(0.5), card_w - Inches(0.6), Inches(1.6),
                      '', font_size=11)
    for item in items:
        add_para(tf, f'·  {item}', font_size=11, color=TEXT_SECONDARY)


# ════════════════════════════════════════
# SLIDE 6: SS500 PM + DVP/QC
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             'SS500 — PM 역할 · 양산 품질 체계', font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3))

# PM card
add_rect(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3.0), BG_CARD, border=True)
add_text_box(slide, Inches(1.1), Inches(1.45), Inches(5), Inches(0.3),
             '제품개발팀 PM 역할', font_size=16, color=GREEN, bold=True)

pm_items = [
    ('이슈 관리', '37건+ 발행/트래킹, 직접 담당 14건'),
    ('BOM · 부품', '60+ 부품 스펙 확정, 참고자료 183건'),
    ('호기 관리', 'EVT 3대 + PVT 16대, 4단계 명명 체계 직접 설계'),
    ('양산 일정', '초도품 16대, 일정 변동 실시간 관리'),
    ('다부서 조율', '전력제어/융합제어/HW/SW/SCM/MD 6개 팀'),
]
y = Inches(1.9)
for label, value in pm_items:
    tf = add_text_box(slide, Inches(1.1), y, Inches(5), Inches(0.35), '', font_size=12)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f'{label}  '
    run.font.size = Pt(11)
    run.font.color.rgb = ACCENT
    run.font.bold = True
    run.font.name = '맑은 고딕'
    run = p.add_run()
    run.text = value
    run.font.size = Pt(12)
    run.font.color.rgb = TEXT_PRIMARY
    run.font.name = '맑은 고딕'
    y += Inches(0.42)

# DVP card
add_rect(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(1.8), BG_CARD, border=True)
add_text_box(slide, Inches(7.1), Inches(1.45), Inches(5), Inches(0.3),
             'DVP — 설계 검증', font_size=14, color=ACCENT, bold=True)

dvp_items = [
    ('DVP.1  구동부 내구성', '완료', GREEN),
    ('DVP.2  자율주행 성능', '진행 중', YELLOW),
    ('DVP.3  펌프 성능', '진행 중', YELLOW),
    ('DVP.4  풍량 성능', '진행 중', YELLOW),
    ('DVP.5  KGAT 종합 검정', '진행 중', YELLOW),
]
y = Inches(1.85)
for item, status, color in dvp_items:
    tf = add_text_box(slide, Inches(7.1), y, Inches(4), Inches(0.25), '', font_size=11)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = item
    run.font.size = Pt(11)
    run.font.color.rgb = TEXT_SECONDARY
    run.font.name = '맑은 고딕'
    run = p.add_run()
    run.text = f'  {status}'
    run.font.size = Pt(10)
    run.font.color.rgb = color
    run.font.bold = True
    run.font.name = '맑은 고딕'
    y += Inches(0.25)

# QC 5단계 card
add_rect(slide, Inches(6.8), Inches(3.4), Inches(5.5), Inches(2.0), BG_CARD, border=True)
add_text_box(slide, Inches(7.1), Inches(3.55), Inches(5), Inches(0.3),
             'QC 5단계 — 양산 품질 체계 구축', font_size=14, color=ACCENT, bold=True)

# QC flow
add_text_box(slide, Inches(7.1), Inches(3.95), Inches(5), Inches(0.3),
             '업체출하 → IQC → LQC → EOL → OQC',
             font_size=12, color=WHITE, bold=True, font_name='Consolas')

tf = add_text_box(slide, Inches(7.1), Inches(4.35), Inches(5), Inches(1), '', font_size=11)
add_para(tf, '전체 36건 검사 문서 중 17건 작성 (47%)', font_size=12, color=TEXT_PRIMARY)
add_para(tf, '업체 출하검사 11건 중 5건 직접 작성', font_size=11, color=TEXT_SECONDARY)
add_para(tf, 'IQC 28개 부품 — 10건+ 가이드 (V0.1~V0.5)', font_size=11, color=TEXT_SECONDARY)

# Tech stack
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.3),
             'STM32F479  ·  CAN 2.0B  ·  CAN FD  ·  BMS  ·  APQP/PPAP  ·  Notion  ·  C',
             font_size=11, color=TEXT_MUTED, font_name='Consolas', alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# SLIDE 7: Research Background (condensed)
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             'Research — PMSM 고장진단', font_size=28, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.35),
             '건국대 RBDO 연구실  |  120kW IPMSM 모델링 및 고장진단  |  현대자동차·현대엔지비 과제',
             font_size=13, color=TEXT_MUTED)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3))

# Research pipeline (simplified)
stages = [
    ('고장 신호\n분석', '모터 및 인버터\n제어 시뮬레이션', '실험적 검증\n(다이나모)', '고장 신호\n검증', '시뮬레이션\n모델 개선'),
]

stage_w = Inches(2.0)
start_x = Inches(0.8)
for i, stage in enumerate(stages[0]):
    x = start_x + i * (stage_w + Inches(0.3))
    is_highlight = i in [1, 2]
    add_rect(slide, x, Inches(1.6), stage_w, Inches(0.8),
             ACCENT_DIM if is_highlight else BG_CARD, border=True)
    add_text_box(slide, x + Inches(0.1), Inches(1.65), stage_w - Inches(0.2), Inches(0.7),
                 stage, font_size=10, color=WHITE if is_highlight else TEXT_SECONDARY,
                 alignment=PP_ALIGN.CENTER, bold=is_highlight)
    if i < 4:
        add_text_box(slide, x + stage_w + Inches(0.05), Inches(1.7), Inches(0.2), Inches(0.5),
                     '→', font_size=14, color=TEXT_MUTED, font_name='Consolas')

# Key contributions in 3 cards
research_cards = [
    ('모터 모델링', [
        '80kW, 120kW급 IPMSM FEA 해석',
        'Flux State Variable Model 적용',
        '포화특성/공간고조파/손실 반영',
        '3D LUT (쇄교자속/토크) 생성',
    ]),
    ('제어 시뮬레이션', [
        'FOC + SVPWM 구현 (Simulink)',
        'MTPA-FW-MTPV 영역 제어',
        '인버터 스위칭 모델 정밀도 비교',
        '시뮬-실험 오차 기반 모델 개선',
    ]),
    ('실험 검증', [
        '200kW 다이나모 테스트 벤드',
        '감자/편심 고장 모터 직접 제작',
        '8가지 운전 조건 시험 설계',
        'DAQ (LabVIEW) 데이터 수집',
    ]),
]

for i, (title, items) in enumerate(research_cards):
    x = Inches(0.8) + i * (Inches(3.8) + Inches(0.3))
    add_rect(slide, x, Inches(2.8), Inches(3.8), Inches(2.5), BG_CARD, border=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.95), Inches(3.2), Inches(0.3),
                 title, font_size=14, color=ACCENT, bold=True)
    tf = add_text_box(slide, x + Inches(0.3), Inches(3.35), Inches(3.2), Inches(1.8), '', font_size=11)
    for item in items:
        add_para(tf, f'·  {item}', font_size=11, color=TEXT_SECONDARY)

# Results
add_rect(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.0), BG_CARD, border=True)
add_text_box(slide, Inches(1.1), Inches(5.7), Inches(10.5), Inches(0.3),
             '연구 성과', font_size=14, color=GREEN, bold=True)
add_text_box(slide, Inches(1.1), Inches(6.05), Inches(10.5), Inches(0.4),
             '감자 및 편심 결함의 이론적 고장진단 신호 도출 → Co-simulation 검증 → 200kW 다이나모 실험 검증  |  시뮬레이션-실험 경향성 일치 확인',
             font_size=12, color=TEXT_SECONDARY)


# ════════════════════════════════════════
# SLIDE 8: Engineering Philosophy
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             'Engineering Philosophy', font_size=32, color=WHITE, bold=True, font_name='Consolas')
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3))

# Philosophy 1: Root Cause
add_rect(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3.5), BG_CARD, border=True)
add_text_box(slide, Inches(1.1), Inches(1.45), Inches(5), Inches(0.3),
             '"정확한 원인 판별이 효율적인 개발의 출발점"',
             font_size=14, color=ACCENT, bold=True)
tf = add_text_box(slide, Inches(1.1), Inches(1.9), Inches(5), Inches(1.2), '', font_size=12)
add_para(tf, '시스템/HW/SW/환경 중 근본 원인을 구분하고', font_size=12, color=TEXT_PRIMARY)
add_para(tf, '올바른 방향에 리소스를 집중합니다.', font_size=12, color=TEXT_PRIMARY)

add_para(tf, '', font_size=6)
add_para(tf, '사례 1: CAN Bus-Off', font_size=12, color=YELLOW, bold=True)
add_para(tf, '통신 에러 → 솔레노이드 GND 바운스 → 배선/접지 문제', font_size=11, color=TEXT_SECONDARY)
add_para(tf, '', font_size=6)
add_para(tf, '사례 2: 극저온 탈조', font_size=12, color=YELLOW, bold=True)
add_para(tf, '제어 문제 → 윤활유 점도↑ → 관로 저항↑ → 기계적 환경', font_size=11, color=TEXT_SECONDARY)

# Philosophy 2: First Principles
add_rect(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(1.8), BG_CARD, border=True)
add_text_box(slide, Inches(7.1), Inches(1.45), Inches(5), Inches(0.3),
             '제1원칙 사고 (First Principles)', font_size=14, color=ACCENT, bold=True)
tf = add_text_box(slide, Inches(7.1), Inches(1.85), Inches(5), Inches(1), '', font_size=12)
add_para(tf, '"관례가 아닌 물리 법칙에서 출발한다"', font_size=12, color=WHITE)
add_para(tf, '', font_size=4)
add_para(tf, '· 호기 관리 체계 부재 → 4단계 명명 체계 직접 설계', font_size=11, color=TEXT_SECONDARY)
add_para(tf, '· CAN 프로토콜 혼재 → Single Source of Truth 확정', font_size=11, color=TEXT_SECONDARY)

# Philosophy 3: 5-step process
add_rect(slide, Inches(6.8), Inches(3.4), Inches(5.5), Inches(2.8), BG_CARD, border=True)
add_text_box(slide, Inches(7.1), Inches(3.55), Inches(5), Inches(0.3),
             '개발 5단계 프로세스', font_size=14, color=ACCENT, bold=True)

steps = [
    ('1', '요구사항을 의심하라', '불필요한 요구사항이 가장 큰 낭비'),
    ('2', '삭제하라', '나중에 추가 > 불필요한 유지'),
    ('3', '단순화/최적화하라', '없어야 할 것을 최적화하지 마라'),
    ('4', '속도를 높여라', '1~3을 먼저 한 다음에만'),
    ('5', '자동화하라', '마지막 단계'),
]
y = Inches(3.95)
for num, title, desc in steps:
    tf = add_text_box(slide, Inches(7.1), y, Inches(5), Inches(0.35), '', font_size=12)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f'{num}. {title}  '
    run.font.size = Pt(12)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = '맑은 고딕'
    run = p.add_run()
    run.text = f'— {desc}'
    run.font.size = Pt(10)
    run.font.color.rgb = TEXT_MUTED
    run.font.name = '맑은 고딕'
    y += Inches(0.38)


# ════════════════════════════════════════
# SLIDE 9: Tech Stack + Publications
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             'Tech Stack & Publications', font_size=28, color=WHITE, bold=True, font_name='Consolas')
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3))

# Tech Stack - HW
add_rect(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3.0), BG_CARD, border=True)
add_text_box(slide, Inches(1.1), Inches(1.45), Inches(5), Inches(0.3),
             'Hardware & Control', font_size=14, color=ACCENT, bold=True)

hw_skills = [
    ('모터 제어', 'FOC, FWC, MTPA, Sensorless', 5),
    ('임베디드', 'STM32, NXP, C/C++', 5),
    ('통신', 'CAN 2.0B, CAN FD, SPI', 5),
    ('전력 시스템', 'BMS, DCDC, 배터리, 릴레이', 4),
    ('시뮬레이션', 'Ansys Maxwell, MATLAB/Simulink', 4),
    ('시험 장비', '다이나모, 오실로스코프, LabVIEW', 4),
]
y = Inches(1.85)
for name, detail, level in hw_skills:
    tf = add_text_box(slide, Inches(1.1), y, Inches(5), Inches(0.3), '', font_size=11)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f'{name}  '
    run.font.size = Pt(11)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = '맑은 고딕'
    run = p.add_run()
    run.text = detail
    run.font.size = Pt(10)
    run.font.color.rgb = TEXT_SECONDARY
    run.font.name = '맑은 고딕'
    run = p.add_run()
    run.text = '  ' + '●' * level + '○' * (5 - level)
    run.font.size = Pt(10)
    run.font.color.rgb = ACCENT
    run.font.name = 'Consolas'
    y += Inches(0.35)

# Tech Stack - SW
add_text_box(slide, Inches(1.1), y + Inches(0.1), Inches(5), Inches(0.3),
             'Software & Tools', font_size=14, color=ACCENT, bold=True)
sw_skills = [
    ('프로그래밍', 'Python, C, JS/TS', 4),
    ('AI 활용', 'Claude Code, AI 워크플로우', 5),
    ('프로젝트 관리', 'Notion, Git, APQP/PPAP', 4),
]
y += Inches(0.5)
for name, detail, level in sw_skills:
    tf = add_text_box(slide, Inches(1.1), y, Inches(5), Inches(0.3), '', font_size=11)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f'{name}  '
    run.font.size = Pt(11)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = '맑은 고딕'
    run = p.add_run()
    run.text = detail
    run.font.size = Pt(10)
    run.font.color.rgb = TEXT_SECONDARY
    run.font.name = '맑은 고딕'
    run = p.add_run()
    run.text = '  ' + '●' * level + '○' * (5 - level)
    run.font.size = Pt(10)
    run.font.color.rgb = ACCENT
    run.font.name = 'Consolas'
    y += Inches(0.35)

# Publications
add_rect(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(5.0), BG_CARD, border=True)
add_text_box(slide, Inches(7.1), Inches(1.45), Inches(5), Inches(0.3),
             'Publications & Patents', font_size=14, color=ACCENT, bold=True)

pubs = [
    ('직무발명', [
        '초기위치추정 알고리즘 (EOP, 2023)',
        '극저온 기동 알고리즘 (EOP, 2024)',
    ]),
    ('주저자 논문', [
        'Co-simulation for Fault Diagnosis of 120kW IPMSM\nand Experimental Validation',
    ]),
    ('공저자 (SCI)', [
        'Condition Monitoring of IGBT Modules using\nInverter Output Parameters\n(IEEE Trans. Reliability)',
    ]),
    ('학회 발표', [
        '한국신뢰성학회 2022 — 최우수발표 논문상',
        '한국PHM학회 2021 — 우수포스터상',
        'PCIM Asia 2022 — IGBT Power Cycling',
    ]),
]
y = Inches(1.85)
for section, items in pubs:
    add_text_box(slide, Inches(7.1), y, Inches(5), Inches(0.25),
                 section, font_size=12, color=GREEN, bold=True)
    y += Inches(0.3)
    for item in items:
        lines = item.split('\n')
        for j, line in enumerate(lines):
            add_text_box(slide, Inches(7.3), y, Inches(4.8), Inches(0.22),
                         line, font_size=10, color=TEXT_SECONDARY if j > 0 else TEXT_PRIMARY)
            y += Inches(0.22)
        y += Inches(0.05)
    y += Inches(0.15)


# ════════════════════════════════════════
# SLIDE 10: Contact / End
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# accent line
add_rect(slide, 0, 0, SLIDE_W, Pt(4), ACCENT)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(0.8),
             'Thank you.', font_size=48, color=WHITE, bold=True, font_name='Consolas',
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.6),
             '황인혁  |  Motor Control & Product Engineer',
             font_size=20, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Contact cards
contacts = [
    ('GitHub', 'github.com/hwanginhyeok'),
    ('Email', 'inhyeok.hwang@gintlab.com'),
    ('Phone', '010-4488-5560'),
]
total_w = len(contacts) * Inches(3.5) + (len(contacts) - 1) * Inches(0.3)
start_x = (Inches(13.333) - total_w) // 2 + Inches(0.5)

for i, (label, value) in enumerate(contacts):
    x = start_x + i * (Inches(3.5) + Inches(0.3))
    add_rect(slide, x, Inches(4.5), Inches(3.5), Inches(1.0), BG_CARD, border=True)
    add_text_box(slide, x + Inches(0.3), Inches(4.6), Inches(3), Inches(0.3),
                 label, font_size=12, color=ACCENT, bold=True, font_name='Consolas')
    add_text_box(slide, x + Inches(0.3), Inches(4.95), Inches(3), Inches(0.3),
                 value, font_size=13, color=TEXT_PRIMARY)


# ── 저장 ──
output_path = os.path.expanduser('~/projects/포트폴리오/포트폴리오_황인혁_2026.pptx')
prs.save(output_path)
print(f'생성 완료: {output_path}')
print(f'슬라이드 수: {len(prs.slides)}')
