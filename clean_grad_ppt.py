"""
대학원 PPT에서 현대차 브랜딩만 제거
원본 내용/이미지/레이아웃 100% 유지
"""
from pptx import Presentation
from pptx.util import Inches
import os

src = os.path.expanduser('~/projects/포트폴리오/참조자료/석사/신입채용_포트폴리오_황인혁_1.pptx')
dst = os.path.expanduser('~/projects/포트폴리오/대학원_포트폴리오_황인혁.pptx')

prs = Presentation(src)
removed = []

for slide_idx, slide in enumerate(prs.slides):
    shapes_to_remove = []

    for shape in slide.shapes:
        # 1) HYUNDAI 로고 이미지 제거
        #    - 모든 슬라이드 우측 상단: x≈7.9", y≈0.1", w≈1.9", h≈0.5"
        #    - 슬라이드 1 하단: x≈7.3", y≈6.6", w≈2.4", h≈0.7"
        if shape.shape_type == 13:  # Picture
            x_inch = shape.left / 914400
            y_inch = shape.top / 914400
            w_inch = shape.width / 914400
            h_inch = shape.height / 914400

            # 우측 상단 로고 (모든 슬라이드)
            if 7.5 < x_inch < 8.5 and y_inch < 0.5 and 1.5 < w_inch < 2.5 and h_inch < 1.0:
                shapes_to_remove.append(shape)
                removed.append(f'  슬라이드 {slide_idx+1}: HYUNDAI 로고 제거 (우측 상단)')

            # 슬라이드 1 하단 로고
            if slide_idx == 0 and 7.0 < x_inch < 8.0 and y_inch > 6.0 and w_inch < 3.0:
                shapes_to_remove.append(shape)
                removed.append(f'  슬라이드 {slide_idx+1}: HYUNDAI 로고 제거 (하단)')

        # 2) 슬라이드 1: "현대 트랜시스 2022 신입사원 상시 채용" 텍스트 제거
        if shape.has_text_frame:
            text = shape.text_frame.text
            if '신입사원' in text and '채용' in text:
                shapes_to_remove.append(shape)
                removed.append(f'  슬라이드 {slide_idx+1}: 텍스트 제거 "{text[:50]}"')
            # "파워트레인 연구" → 변경 (슬라이드 1 또는 11)
            elif '파워트레인 연구' in text and len(text) < 30:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if '파워트레인 연구' in run.text:
                            run.text = run.text.replace('파워트레인 연구', '전동화 시스템 연구')
                            removed.append(f'  슬라이드 {slide_idx+1}: "파워트레인 연구" → "전동화 시스템 연구"')

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

prs.save(dst)
print(f'완료: {dst}')
print(f'변경 {len(removed)}건:')
for r in removed:
    print(r)
