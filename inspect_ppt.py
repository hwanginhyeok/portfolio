"""원본 PPT의 모든 shape 분석 - 현대차 로고 위치 파악"""
from pptx import Presentation
from pptx.util import Inches, Emu
import os

src = os.path.expanduser('~/projects/포트폴리오/참조자료/석사/신입채용_포트폴리오_황인혁_1.pptx')
prs = Presentation(src)

for i, slide in enumerate(prs.slides):
    print(f'\n=== 슬라이드 {i+1} ===')
    for shape in slide.shapes:
        info = f'  type={shape.shape_type} name="{shape.name}" '
        info += f'pos=({shape.left},{shape.top}) size=({shape.width},{shape.height}) '
        info += f'[x={round(shape.left/914400,1)}" y={round(shape.top/914400,1)}" w={round(shape.width/914400,1)}" h={round(shape.height/914400,1)}"]'
        if shape.has_text_frame:
            text = shape.text_frame.text[:60].replace('\n', ' ')
            info += f' text="{text}"'
        if shape.shape_type == 13:  # Picture
            info += ' [IMAGE]'
        print(info)
