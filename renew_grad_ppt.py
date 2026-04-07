"""
대학원 포트폴리오 PPT 디자인 리뉴얼
- 기존 내용 유지, 현대차 브랜딩 제거
- 흰색 깔끔한 테마로 변경
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 디자인 시스템 (화이트 클린) ──
BG = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)       # slate-50
BG_CARD = RGBColor(0xF1, 0xF5, 0xF9)        # slate-100
ACCENT = RGBColor(0x1E, 0x40, 0xAF)         # blue-800
ACCENT_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)   # blue-100
TEXT_PRIMARY = RGBColor(0x0F, 0x17, 0x2A)    # slate-900
TEXT_SECONDARY = RGBColor(0x47, 0x55, 0x69)  # slate-600
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)      # slate-400
BORDER = RGBColor(0xE2, 0xE8, 0xF0)         # slate-200
GREEN = RGBColor(0x16, 0xA3, 0x4A)          # green-600
RED = RGBColor(0xDC, 0x26, 0x26)            # red-600

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=14,
             color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT, font='맑은 고딕'):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf

def add_para(tf, text, size=14, color=TEXT_PRIMARY, bold=False, font='맑은 고딕',
             space_before=Pt(3), alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.space_before = space_before
    p.alignment = alignment
    return p

def add_rect(slide, left, top, width, height, fill=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.02
    return shape

def add_line(slide, left, top, width, color=ACCENT, height=Pt(3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def header_bar(slide, title, subtitle=None):
    """페이지 상단 헤더"""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), ACCENT)
    add_text(slide, Inches(0.8), Inches(0.15), Inches(8), Inches(0.6),
             title, size=26, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(0.55), Inches(8), Inches(0.3),
                 subtitle, size=12, color=RGBColor(0xBF, 0xDB, 0xFE))

def page_num(slide, num):
    add_text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
             str(num), size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT, font='Consolas')

def research_pipeline(slide, highlight_idx):
    """연구 파이프라인 5단계 (반복 사용)"""
    stages = ['고장 신호\n분석', '모터 및 인버터\n제어 시뮬레이션', '실험적 검증\n(다이나모 테스트)',
              '고장 신호\n검증 (결과)', '시뮬레이션\n모델 개선']
    stage_w = Inches(2.1)
    start_x = Inches(0.8)
    y = Inches(1.2)
    for i, stage in enumerate(stages):
        x = start_x + i * (stage_w + Inches(0.2))
        is_hl = (i == highlight_idx)
        fill = ACCENT if is_hl else BG_CARD
        text_c = RGBColor(0xFF, 0xFF, 0xFF) if is_hl else TEXT_SECONDARY
        add_rect(slide, x, y, stage_w, Inches(0.7), fill, BORDER if not is_hl else None)
        add_text(slide, x + Inches(0.1), y + Inches(0.05), stage_w - Inches(0.2), Inches(0.6),
                 stage, size=10, color=text_c, align=PP_ALIGN.CENTER, bold=is_hl)
        if i < 4:
            add_text(slide, x + stage_w + Inches(0.02), y + Inches(0.1), Inches(0.15), Inches(0.4),
                     '→', size=14, color=TEXT_MUTED, font='Consolas')


# ══════════════════════════════════════════
# SLIDE 1: 표지
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_rect(slide, 0, 0, SLIDE_W, Inches(3.5), ACCENT)

add_text(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8),
         '포트폴리오', size=18, color=RGBColor(0xBF, 0xDB, 0xFE))
add_text(slide, Inches(1.5), Inches(1.3), Inches(10), Inches(1.0),
         '전동화 시스템 버추얼 개발', size=40, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
add_text(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(0.5),
         'PMSM 고장진단을 위한 Co-simulation 및 실험 검증',
         size=16, color=RGBColor(0xBF, 0xDB, 0xFE))

add_text(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.6),
         '황인혁', size=32, color=TEXT_PRIMARY, bold=True)
add_text(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.4),
         '건국대학교 기계설계학과 석사  |  RBDO 연구실 (신뢰성기반최적설계)',
         size=14, color=TEXT_SECONDARY)

# bottom accent line
add_line(slide, 0, Inches(7.3), SLIDE_W, ACCENT, Pt(4))
add_text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.4),
         '"모터 인버터 시스템 연구 경험, 최적을 모델을 찾는 끈기"',
         size=14, color=ACCENT, bold=True, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════
# SLIDE 2: 이력서
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_bar(slide, '이력서')
page_num(slide, 1)

# Profile section
add_text(slide, Inches(2.5), Inches(1.2), Inches(8), Inches(0.35),
         '이 름 : 황인혁', size=14, color=TEXT_PRIMARY, bold=True)
profile_lines = [
    '학 력 : 건국대학교 기계공학과 학사 (\'21. 2월 졸업)',
    '           건국대학교 기계설계학과 석사 (\'23. 2월 졸업)',
    'LAB실명 : 건국대학교 신뢰성기반최적설계 연구실',
    '세부전공 : 매입형 영구자석 동기전동기의 모델링 및 고장 진단에 관한 연구',
]
y = Inches(1.6)
for line in profile_lines:
    add_text(slide, Inches(2.5), y, Inches(8), Inches(0.3), line, size=12, color=TEXT_SECONDARY)
    y += Inches(0.3)

# 프로필 사진 자리 (placeholder)
add_rect(slide, Inches(1.0), Inches(1.2), Inches(1.3), Inches(1.7), BG_CARD, BORDER)
add_text(slide, Inches(1.0), Inches(1.8), Inches(1.3), Inches(0.3),
         '사진', size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# 프로젝트 경력
add_rect(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(3.5), BG_LIGHT, BORDER)
add_rect(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(0.4), ACCENT)
add_text(slide, Inches(1.0), Inches(3.22), Inches(5), Inches(0.35),
         '프로젝트 경력', size=13, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)

tf = add_text(slide, Inches(1.0), Inches(3.8), Inches(5), Inches(2.8), '', size=11)
add_para(tf, '• 2년 6개월 연구 경력 (학부: 6개월 / 석사: 2년)', size=11, color=TEXT_PRIMARY, bold=True)
add_para(tf, '', size=6)
add_para(tf, '• 주요 프로젝트', size=11, color=TEXT_PRIMARY, bold=True)
add_para(tf, '  - 전동화차량 핵심부품의 고장예지 기술개발 2단계', size=10, color=TEXT_SECONDARY)
add_para(tf, '    (2020-11 ~ 2022-07, 현대자동차 \'상용 해석 리서치랩\')', size=10, color=TEXT_MUTED)
add_para(tf, '  - 전동화 차량 구동 전기모터의 상태 진단을 위한 PHM SoC 개발', size=10, color=TEXT_SECONDARY)
add_para(tf, '    (2021-04 ~ 2022-10, 산업통상자원부)', size=10, color=TEXT_MUTED)
add_para(tf, '  - 스마트 차량 정비를 위한 핵심 부품의 잔류수명 예측기술 및', size=10, color=TEXT_SECONDARY)
add_para(tf, '    비접촉 센서 모듈 개발 (2021-07 ~ 2022-10, 만도)', size=10, color=TEXT_MUTED)
add_para(tf, '', size=6)
add_para(tf, '• 대외활동 경험', size=11, color=TEXT_PRIMARY, bold=True)
add_para(tf, '  - 교내봉사 동아리 MRA (2015-09 ~ 2016-09)', size=10, color=TEXT_SECONDARY)
add_para(tf, '  - 미식축구 동아리 Raging Bulls (2015-07 ~ 2019-08)', size=10, color=TEXT_SECONDARY)
add_para(tf, '  - 수상레저 동아리 SUPREME (2020-05 ~ 2021-10)', size=10, color=TEXT_SECONDARY)

# 주요 발표 논문 및 연구실적
add_rect(slide, Inches(6.8), Inches(3.2), Inches(5.5), Inches(3.5), BG_LIGHT, BORDER)
add_rect(slide, Inches(6.8), Inches(3.2), Inches(5.5), Inches(0.4), ACCENT)
add_text(slide, Inches(7.0), Inches(3.22), Inches(5), Inches(0.35),
         '주요 발표 논문 및 연구실적', size=13, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)

tf = add_text(slide, Inches(7.0), Inches(3.8), Inches(5), Inches(2.8), '', size=11)
add_para(tf, '• 주저자 논문', size=11, color=TEXT_PRIMARY, bold=True)
add_para(tf, '  - Co-simulation for fault diagnosis of 120kW Interior Permanent', size=10, color=TEXT_SECONDARY)
add_para(tf, '    Magnet Synchronous Machine and Experimental Validation', size=10, color=TEXT_SECONDARY)
add_para(tf, '', size=6)
add_para(tf, '• 주요 학회 발표 실적', size=11, color=TEXT_PRIMARY, bold=True)
add_para(tf, '  - 전기자동차용 매입형 영구자석 동기전동기의 정밀한 시뮬레이션을 위한', size=10, color=TEXT_SECONDARY)
add_para(tf, '    모델링 분석 (한국 PHM 학회 정기학술대회 2022)', size=10, color=TEXT_SECONDARY)
add_para(tf, '  - 전동화 차량 구동시스템의 효율적인 예방정비 기술 개발', size=10, color=TEXT_SECONDARY)
add_para(tf, '    (한국신뢰성학회 2022년 춘계학술대회, 최우수발표 논문상)', size=10, color=ACCENT, bold=True)
add_para(tf, '  - 시스템 수준 측정값을 이용한 모터 구동 시스템 내 IGBT 개방 고장', size=10, color=TEXT_SECONDARY)
add_para(tf, '    진단 기법 (한국 PHM 학회 2021, 우수포스터상)', size=10, color=ACCENT, bold=True)


# ══════════════════════════════════════════
# SLIDE 3: 이수 전공 Positioning Map
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_bar(slide, '이수 전공 Positioning Map')
page_num(slide, 2)

# 학사/석사 학점
tf = add_text(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.8), '', size=12)
add_para(tf, '학사과정', size=14, color=TEXT_PRIMARY, bold=True)
add_para(tf, '- 전공이수학점: 79학점 이수 (3.24/4.5)', size=11, color=TEXT_SECONDARY)
add_para(tf, '- 3학년 3.11 → 4학년 3.58', size=11, color=TEXT_SECONDARY)

tf = add_text(slide, Inches(6.5), Inches(1.2), Inches(5), Inches(0.8), '', size=12)
add_para(tf, '석사과정', size=14, color=TEXT_PRIMARY, bold=True)
add_para(tf, '- 전공이수학점: 27학점 이수 (3.77/4.5)', size=11, color=TEXT_SECONDARY)
add_para(tf, '- 스위칭제어전원, 신뢰성기반최적설계', size=11, color=TEXT_SECONDARY)

# Positioning Map 가이드
add_text(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(0.3),
         '※ 아래 Positioning Map은 원본 PPT의 과목-관심도-성적 매트릭스를 참조하여 노트북에서 이미지로 삽입하세요.',
         size=11, color=RED)

# 직무 관련 과목 하이라이트
add_rect(slide, Inches(0.8), Inches(2.9), Inches(5.5), Inches(2.0), BG_LIGHT, BORDER)
add_text(slide, Inches(1.0), Inches(3.0), Inches(5), Inches(0.3),
         '관심도가 높았던 과목 (설계과목)', size=13, color=ACCENT, bold=True)
tf = add_text(slide, Inches(1.0), Inches(3.4), Inches(5), Inches(1.3), '', size=11)
add_para(tf, '- 창의설계 : 코딩을 사용하여 목적에 맞는 로봇 구현 (A+)', size=11, color=TEXT_SECONDARY)
add_para(tf, '- 종합설계 : 전동화 시스템 모델링 및 인버터 개방회로 시뮬레이션 (A+)', size=11, color=TEXT_SECONDARY)
add_para(tf, '- 신뢰성기반최적설계 : 시험과 신뢰성 이론에 이해도 (A)', size=11, color=TEXT_SECONDARY)

add_rect(slide, Inches(6.8), Inches(2.9), Inches(5.5), Inches(2.0), BG_LIGHT, BORDER)
add_text(slide, Inches(7.0), Inches(3.0), Inches(5), Inches(0.3),
         'C/C+ 수강 과목', size=13, color=TEXT_MUTED, bold=True)
tf = add_text(slide, Inches(7.0), Inches(3.4), Inches(5), Inches(1.3), '', size=11)
add_para(tf, '- 학부 전공 : 동아리와 아르바이트로 인한 시간관리 미흡', size=11, color=TEXT_SECONDARY)
add_para(tf, '- 석사 전공 : 타과 전공 지식 및 툴(ROS, HETLAS) 사용능력 부족', size=11, color=TEXT_SECONDARY)
add_para(tf, '- 석사생활 동안 학습과 연구를 통한 개선', size=11, color=TEXT_SECONDARY)

# 직무관련 전공 범례
legend_items = [('학사 타전공', BG_CARD), ('학사 전공', ACCENT_LIGHT),
                ('석사 전공', RGBColor(0xBF, 0xDB, 0xFE)), ('직무관련전공', ACCENT)]
x = Inches(0.8)
for label, color in legend_items:
    add_rect(slide, x, Inches(5.3), Inches(1.2), Inches(0.3), color, BORDER)
    add_text(slide, x + Inches(0.1), Inches(5.3), Inches(1.0), Inches(0.3),
             label, size=9, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)
    x += Inches(1.5)

# 주요 전공 (높은 관심 + 좋은 성적)
add_rect(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), BG_LIGHT, BORDER)
add_text(slide, Inches(1.0), Inches(5.9), Inches(10), Inches(0.3),
         '직무 핵심 과목 (높은 관심 + A/A+ 성적)', size=12, color=ACCENT, bold=True)
subjects = ['제어공학', '스위칭제어전원', '계측공학', '신뢰성기반최적설계',
            '신뢰성공학', '종합설계', '창의설계', '마이크로컨트롤러설계및응용']
tf = add_text(slide, Inches(1.0), Inches(6.3), Inches(10), Inches(0.5), '', size=11)
add_para(tf, '  ·  '.join(subjects), size=11, color=TEXT_PRIMARY)


# ══════════════════════════════════════════
# SLIDE 4: 주요 연구 — PMSM 고장진단 개요
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_bar(slide, '주요 연구 수행 이력', 'PMSM의 고장진단')
page_num(slide, 3)
research_pipeline(slide, -1)  # overview, no highlight

# 제목
add_text(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.4),
         'PMSM의 감자 및 편심 결함 고장진단에 관한 연구', size=18, color=TEXT_PRIMARY, bold=True)

# 연구배경
add_rect(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.2), BG_LIGHT, BORDER)
add_text(slide, Inches(1.0), Inches(2.85), Inches(1.5), Inches(0.3),
         '연구배경', size=13, color=ACCENT, bold=True)
tf = add_text(slide, Inches(2.5), Inches(2.85), Inches(9), Inches(1.0), '', size=11)
add_para(tf, '• 모터의 감자 (demagnetization) 및 편심 (eccentricity) 결함 발생 → 진동 및 소음 / 모터 효율 저감', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 초기 진단 실패 시 unbalanced magnetic pull (UMP) 지속 → 베어링 고장 / 편심', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 상태 모니터링을 통한 상태 기반 유지보수 필요 → 모터의 고장진단 기술', size=11, color=TEXT_SECONDARY)

# 연구내용
add_rect(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.5), BG_LIGHT, BORDER)
add_text(slide, Inches(1.0), Inches(4.35), Inches(1.5), Inches(0.3),
         '연구내용', size=13, color=ACCENT, bold=True)
tf = add_text(slide, Inches(2.5), Inches(4.35), Inches(9), Inches(2.2), '', size=11)
add_para(tf, '1) 모터 영상과 제어기 영향을 고려한 고장신호 선정', size=11, color=TEXT_PRIMARY, bold=True)
add_para(tf, '   감자 고장 발생시 전류 증가, 편심 고장 발생시 전류 fractional 고조파 증가', size=10, color=TEXT_SECONDARY)
add_para(tf, '2) 고장 모드를 반영한 모터의 전자기 해석 수행 (Ansys Maxwell)', size=11, color=TEXT_PRIMARY, bold=True)
add_para(tf, '   감자 고장 쇄교자속 감소, 편심 고장 쇄교자속 왜곡', size=10, color=TEXT_SECONDARY)
add_para(tf, '3) 시뮬레이션을 통한 검증 : 정확성과 해석시간을 고려한 모터 및 인버터 모델링', size=11, color=TEXT_PRIMARY, bold=True)
add_para(tf, '4) 실험적 검증을 위한 정상 및 고장 모터 제작, 200kW급 다이나모 테스트', size=11, color=TEXT_PRIMARY, bold=True)
add_para(tf, '5) 시뮬레이션 모델 개선', size=11, color=TEXT_PRIMARY, bold=True)

# 결론
add_rect(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.35), ACCENT)
add_text(slide, Inches(1.0), Inches(7.0), Inches(11), Inches(0.3),
         '감자 및 편심 결함을 위한 이론적 고장진단 신호 파악 및 시뮬레이션과 실험을 통한 검증',
         size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════
# SLIDE 5: 고장신호 분석
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_bar(slide, '주요 연구 수행 이력', 'PMSM 고장진단 — 1) 모터의 고장신호 분석')
page_num(slide, 4)
research_pipeline(slide, 0)

add_text(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.3),
         '1) 모터의 고장신호 분석', size=16, color=TEXT_PRIMARY, bold=True)
tf = add_text(slide, Inches(0.8), Inches(2.6), Inches(11), Inches(0.6), '', size=11)
add_para(tf, '• 모터 고장 발생시 전류 신호의 변화를 이용한 모터 고장진단 가능', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 고장 발생시 전류의 변화 형태를 사전에 알고 진단하는 것이 중요', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 고장신호는 모터 고장의 종류와 형상에 영향을 받음', size=11, color=TEXT_SECONDARY)

# 두 카드
add_rect(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(3.2), BG_LIGHT, BORDER)
add_rect(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(0.35), ACCENT)
add_text(slide, Inches(1.0), Inches(3.52), Inches(5), Inches(0.3),
         '1. 고장 분석', size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
tf = add_text(slide, Inches(1.0), Inches(4.1), Inches(5), Inches(2.0), '', size=11)
add_para(tf, '• 정상 및 고장 모터에 대한 전자기 해석 진행', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 모터 파라미터 (쇄교자속 및 토크) 추출', size=11, color=TEXT_SECONDARY)
add_para(tf, '', size=8)
add_para(tf, '※ 이미지: 정상 모터 / 감자 모터 FEA 결과', size=10, color=RED)
add_para(tf, '  (원본 PPT에서 이미지 복사-붙여넣기)', size=10, color=RED)

add_rect(slide, Inches(6.8), Inches(3.5), Inches(5.5), Inches(3.2), BG_LIGHT, BORDER)
add_rect(slide, Inches(6.8), Inches(3.5), Inches(5.5), Inches(0.35), ACCENT)
add_text(slide, Inches(7.0), Inches(3.52), Inches(5), Inches(0.3),
         '2. 고장신호 특성화', size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
tf = add_text(slide, Inches(7.0), Inches(4.1), Inches(5), Inches(2.0), '', size=11)
add_para(tf, '• 정상 및 고장 모터에 대한 전자기 해석 진행', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 모터의 쇄교자속 추출 및 신호 분석', size=11, color=TEXT_SECONDARY)
add_para(tf, '', size=8)
add_para(tf, '감자 고장 : 자속 감소', size=11, color=ACCENT, bold=True)
add_para(tf, '복합 편심 : 주기성 (동적 및 정적 편심 자속 증가)', size=11, color=ACCENT, bold=True)

add_rect(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.35), ACCENT_LIGHT)
add_text(slide, Inches(1.0), Inches(6.9), Inches(11), Inches(0.3),
         '감자(전류 증가) 및 편심(fractional harmonics) 결함 신호 이론적 도출',
         size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════
# SLIDE 6: 모터 및 인버터 시뮬레이션 (모델 선정 + 모델링)
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_bar(slide, '주요 연구 수행 이력', 'PMSM 고장진단 — 2) 모터 및 인버터 제어 시뮬레이션')
page_num(slide, 5)
research_pipeline(slide, 1)

add_text(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.3),
         '2) 모터 및 인버터 제어 시뮬레이션', size=16, color=TEXT_PRIMARY, bold=True)

# 모터 모델 선정
add_rect(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(2.0), BG_LIGHT, BORDER)
add_rect(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(0.35), ACCENT)
add_text(slide, Inches(1.0), Inches(2.72), Inches(5), Inches(0.3),
         '1. 모터 모델 선정', size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
tf = add_text(slide, Inches(1.0), Inches(3.2), Inches(5), Inches(1.3), '', size=11)
add_para(tf, '• 모터 결함 반영 → 모델 정확도 중요', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 다양한 동작환경 → 계산시간 중요', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 정확성/계산시간을 모두 고려한 모델 선정', size=11, color=TEXT_SECONDARY)
add_para(tf, '', size=4)
add_para(tf, '→ 포화특성, 공간고조파, 손실을 반영한', size=11, color=ACCENT, bold=True)
add_para(tf, '   Flux State Variable Model 사용', size=11, color=ACCENT, bold=True)

# 모터 모델링
add_rect(slide, Inches(6.8), Inches(2.7), Inches(5.5), Inches(2.0), BG_LIGHT, BORDER)
add_rect(slide, Inches(6.8), Inches(2.7), Inches(5.5), Inches(0.35), ACCENT)
add_text(slide, Inches(7.0), Inches(2.72), Inches(5), Inches(0.3),
         '2. 모터 모델링', size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
tf = add_text(slide, Inches(7.0), Inches(3.2), Inches(5), Inches(1.3), '', size=11)
add_para(tf, '• Transient FEM 해석 (Ansys Maxwell)', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 정상 및 고장 모터에 대한 쇄교자속 / 토크 추출', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 모터 파라미터 → Look-up table (LUT)', size=11, color=TEXT_SECONDARY)
add_para(tf, '  MATLAB Simulink 모터 모델 구현시 활용', size=11, color=TEXT_SECONDARY)
add_para(tf, '', size=4)
add_para(tf, '→ 쇄교자속 / 모터 토크 → 3D LUT (id, iq map)', size=11, color=ACCENT, bold=True)

# 인버터 스위치 모델 + 제어시스템
add_rect(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(2.0), BG_LIGHT, BORDER)
add_rect(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(0.35), ACCENT)
add_text(slide, Inches(1.0), Inches(5.02), Inches(5), Inches(0.3),
         '3. 인버터 스위치 모델 선정', size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
tf = add_text(slide, Inches(1.0), Inches(5.5), Inches(5), Inches(1.3), '', size=11)
add_para(tf, '• 다양한 동작환경 → 계산시간 중요', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 해석시간과 모터의 고장신호의 관찰 가능성 고려', size=11, color=TEXT_SECONDARY)
add_para(tf, '', size=4)
add_para(tf, '→ Turn-on & off / Dead time 반영 모델 선정', size=11, color=ACCENT, bold=True)

add_rect(slide, Inches(6.8), Inches(5.0), Inches(5.5), Inches(2.0), BG_LIGHT, BORDER)
add_rect(slide, Inches(6.8), Inches(5.0), Inches(5.5), Inches(0.35), ACCENT)
add_text(slide, Inches(7.0), Inches(5.02), Inches(5), Inches(0.3),
         '4. 모터 및 제어시스템 해석', size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
tf = add_text(slide, Inches(7.0), Inches(5.5), Inches(5), Inches(1.3), '', size=11)
add_para(tf, '• MATLAB Simulink 이용 모터 모델 구현', size=11, color=TEXT_SECONDARY)
add_para(tf, '• PMSM 속도 제어 시스템 구현', size=11, color=TEXT_SECONDARY)
add_para(tf, '  - SVPWM, FOC, PI current/speed control', size=10, color=TEXT_MUTED)
add_para(tf, '  - MTPA, Field-weakening (FW), MTPV', size=10, color=TEXT_MUTED)

add_rect(slide, Inches(0.8), Inches(7.1), Inches(11.5), Inches(0.3), ACCENT_LIGHT)
add_text(slide, Inches(1.0), Inches(7.1), Inches(11), Inches(0.25),
         '모델의 정확도 / 계산시간을 고려한 정상 및 고장 모터 제어 시스템 모델링',
         size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════
# SLIDE 7: 실험적 검증
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_bar(slide, '주요 연구 수행 이력', 'PMSM 고장진단 — 3) 실험적 검증')
page_num(slide, 6)
research_pipeline(slide, 2)

add_text(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.3),
         '3) 실험적 검증', size=16, color=TEXT_PRIMARY, bold=True)

# 모터 제작
add_rect(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(2.2), BG_LIGHT, BORDER)
add_rect(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(0.35), ACCENT)
add_text(slide, Inches(1.0), Inches(2.72), Inches(5), Inches(0.3),
         '1. 모터 제작', size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
tf = add_text(slide, Inches(1.0), Inches(3.2), Inches(5), Inches(1.5), '', size=11)
add_para(tf, '• 정상 및 고장 모터 제작', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 인위적 고장 삽입', size=11, color=TEXT_SECONDARY)
add_para(tf, '  - 전체 감자 모터 (하우징 + 로터 축)', size=10, color=TEXT_MUTED)
add_para(tf, '  - 정적 & 동적 편심 고장 모터', size=10, color=TEXT_MUTED)
add_para(tf, '', size=6)
add_para(tf, '※ 이미지: 고장 모터 사진 (원본 PPT 참조)', size=10, color=RED)

# 다이나모 테스트
add_rect(slide, Inches(6.8), Inches(2.7), Inches(5.5), Inches(2.2), BG_LIGHT, BORDER)
add_rect(slide, Inches(6.8), Inches(2.7), Inches(5.5), Inches(0.35), ACCENT)
add_text(slide, Inches(7.0), Inches(2.72), Inches(5), Inches(0.3),
         '2. 다이나모 테스트', size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
tf = add_text(slide, Inches(7.0), Inches(3.2), Inches(5), Inches(1.5), '', size=11)
add_para(tf, '• 200kW 다이나모 테스트 벤드 이용', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 기본구성 : 파워스토리지, 모터, 인버터, 다이나모미터', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 추가구성 : 가속도계, 전류센서, 전압센서, DAQ(LabVIEW)', size=11, color=TEXT_SECONDARY)
add_para(tf, '• MTPA-FW-MTPV 영역을 고려한 8가지 실험조건 설계', size=11, color=TEXT_SECONDARY)
add_para(tf, '  (2000/5500/8500rpm × 50/100/250/400Nm)', size=10, color=TEXT_MUTED)
add_para(tf, '', size=6)
add_para(tf, '※ 이미지: 실험 구성도 + 셋업 사진 (원본 PPT)', size=10, color=RED)


# ══════════════════════════════════════════
# SLIDE 8: 고장신호 검증
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_bar(slide, '주요 연구 수행 이력', 'PMSM 고장진단 — 4) 고장신호 검증')
page_num(slide, 7)
research_pipeline(slide, 3)

add_text(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.3),
         '4) 고장신호 검증', size=16, color=TEXT_PRIMARY, bold=True)

add_rect(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(3.6), BG_LIGHT, BORDER)
add_rect(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(0.35), ACCENT)
add_text(slide, Inches(1.0), Inches(2.72), Inches(5), Inches(0.3),
         '1. 고장 경향성 비교 분석', size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
tf = add_text(slide, Inches(1.0), Inches(3.2), Inches(5), Inches(2.8), '', size=11)
add_para(tf, '시뮬레이션 결과와 실험 결과의 경향성 비교', size=11, color=TEXT_PRIMARY, bold=True)
add_para(tf, '', size=6)
add_para(tf, '• Demagnetization: 시뮬-실험 전류 파형 경향 일치', size=11, color=TEXT_SECONDARY)
add_para(tf, '• Mixed Eccentricity: 고조파 스펙트럼 경향 일치', size=11, color=TEXT_SECONDARY)
add_para(tf, '', size=8)
add_para(tf, '※ 이미지: Phase A current 비교 그래프', size=10, color=RED)
add_para(tf, '  + Current spectrum 비교 그래프 (원본 PPT)', size=10, color=RED)

add_rect(slide, Inches(6.8), Inches(2.7), Inches(5.5), Inches(3.6), BG_LIGHT, BORDER)
add_rect(slide, Inches(6.8), Inches(2.7), Inches(5.5), Inches(0.35), ACCENT)
add_text(slide, Inches(7.0), Inches(2.72), Inches(5), Inches(0.3),
         '2. 전류 크기 비교 분석', size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
tf = add_text(slide, Inches(7.0), Inches(3.2), Inches(5), Inches(2.8), '', size=11)
add_para(tf, '시뮬레이션과 실험 결과 비교', size=11, color=TEXT_PRIMARY, bold=True)
add_para(tf, '', size=6)
add_para(tf, '• CSVM: THD= 2.2%', size=11, color=TEXT_SECONDARY)
add_para(tf, '• FSVM: THD= 4.3%', size=11, color=TEXT_SECONDARY)
add_para(tf, '• Experiment: THD= 6.1%', size=11, color=TEXT_SECONDARY)
add_para(tf, '', size=6)
add_para(tf, '→ 시뮬레이션과 실험의 전류 경향성 일치 확인', size=11, color=ACCENT, bold=True)
add_para(tf, '', size=8)
add_para(tf, '※ 이미지: FFT comparison + Phase A current', size=10, color=RED)
add_para(tf, '  비교 그래프 (원본 PPT)', size=10, color=RED)


# ══════════════════════════════════════════
# SLIDE 9: 시뮬레이션 모델 개선
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_bar(slide, '주요 연구 수행 이력', 'PMSM 고장진단 — 5) 시뮬레이션 모델 개선')
page_num(slide, 8)
research_pipeline(slide, 4)

add_text(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.3),
         '5) 시뮬레이션 모델 개선', size=16, color=TEXT_PRIMARY, bold=True)

add_rect(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(3.6), BG_LIGHT, BORDER)
add_rect(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(0.35), ACCENT)
add_text(slide, Inches(1.0), Inches(2.72), Inches(5), Inches(0.3),
         '1. 모델링 개선점 발견', size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
tf = add_text(slide, Inches(1.0), Inches(3.2), Inches(5), Inches(2.8), '', size=11)
add_para(tf, '• 시뮬레이션과 실험 지령 전류각과 크기가 다름', size=11, color=TEXT_SECONDARY)
add_para(tf, '• 전류 제어기의 알고리즘 개선 필요', size=11, color=TEXT_SECONDARY)
add_para(tf, '', size=6)
add_para(tf, '기존의 전류 제어기 전류 지령 맵:', size=11, color=TEXT_PRIMARY, bold=True)
add_para(tf, '시뮬레이션과 실험의 id-iq 전류 지령이 불일치', size=10, color=TEXT_MUTED)
add_para(tf, '', size=8)
add_para(tf, '※ 이미지: 기존 전류 제어기 전류 지령 맵', size=10, color=RED)
add_para(tf, '  (원본 PPT에서 복사)', size=10, color=RED)

add_rect(slide, Inches(6.8), Inches(2.7), Inches(5.5), Inches(3.6), BG_LIGHT, BORDER)
add_rect(slide, Inches(6.8), Inches(2.7), Inches(5.5), Inches(0.35), ACCENT)
add_text(slide, Inches(7.0), Inches(2.72), Inches(5), Inches(0.3),
         '2. 전류 제어기 재설계', size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
tf = add_text(slide, Inches(7.0), Inches(3.2), Inches(5), Inches(2.8), '', size=11)
add_para(tf, '• 실제 전류 제어기는 수학적 모델과 모터 특성을 반영한', size=11, color=TEXT_SECONDARY)
add_para(tf, '  해석 결과를 활용하여 calibration을 하여 설계함', size=11, color=TEXT_SECONDARY)
add_para(tf, '• ANSYS Maxwell 모터 해석결과를 활용하여', size=11, color=TEXT_SECONDARY)
add_para(tf, '  MATLAB CAGE에서 최적화 진행', size=11, color=TEXT_SECONDARY)
add_para(tf, '', size=6)
add_para(tf, '→ 정확한 모델링을 위한 전류 제어기 재설계 및 적용', size=11, color=ACCENT, bold=True)
add_para(tf, '', size=8)
add_para(tf, '※ 이미지: 개선된 전류 제어기 전류 지령 맵', size=10, color=RED)
add_para(tf, '  + id/iq current map (원본 PPT)', size=10, color=RED)

add_rect(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.35), ACCENT_LIGHT)
add_text(slide, Inches(1.0), Inches(6.5), Inches(11), Inches(0.3),
         '정확한 모델링을 위한 전류 제어기 재설계 및 적용',
         size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════
# SLIDE 10: 마무리
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
page_num(slide, 9)

add_rect(slide, 0, Inches(2.5), SLIDE_W, Inches(2.5), ACCENT)

add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.5),
         '황인혁', size=36, color=TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
         'Thank you.', size=40, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
         align=PP_ALIGN.CENTER, font='Consolas')

add_text(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.5),
         '건국대학교 기계설계학과 석사  |  RBDO 연구실',
         size=14, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

# contact
tf = add_text(slide, Inches(3), Inches(5.5), Inches(7), Inches(1.2), '', size=13)
add_para(tf, 'GitHub    github.com/hwanginhyeok', size=13, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
add_para(tf, 'Email      inhyeok.hwang@gintlab.com', size=13, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_line(slide, 0, Inches(7.3), SLIDE_W, ACCENT, Pt(4))


# ── 저장 ──
output_path = os.path.expanduser('~/projects/포트폴리오/대학원_포트폴리오_황인혁_리뉴얼.pptx')
prs.save(output_path)
print(f'생성 완료: {output_path}')
print(f'슬라이드 수: {len(prs.slides)}')
print()
print('※ 이미지가 포함된 슬라이드(4~9)는 원본 PPT에서 이미지를 복사하여 붙여넣어야 합니다.')
print('   원본: 참조자료/석사/신입채용_포트폴리오_황인혁_1.pptx')
